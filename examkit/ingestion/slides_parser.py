"""
Slides parser for PPTX and PDF formats.
"""

import logging
from pathlib import Path
from typing import Any, Dict, List

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation

from examkit.utils.io_utils import ensure_dir


def parse_pptx(path: Path, cache_dir: Path, logger: logging.Logger) -> List[Dict[str, Any]]:
    """
    Extract structured slide information from a PPTX file.
    
    Parameters:
        path (Path): Path to the source PPTX file.
        cache_dir (Path): Directory used to store slide-related cache (e.g., generated image files).
        logger (logging.Logger): Logger used for progress and warning messages.
    
    Returns:
        List[Dict[str, Any]]: A list of slide dictionaries with the following keys:
            - source (str): Fixed value "slides".
            - type (str): Fixed value "pptx".
            - slide_number (int): 1-based slide index.
            - title (str): Slide title text if present, otherwise empty string.
            - content (List[str]): Text blocks from the slide excluding the title.
            - notes (str): Slide notes text if present, otherwise empty string.
            - images (List[str]): Filenames (placeholders) for images detected on the slide.
    """
    logger.info(f"Parsing PPTX: {path}")

    prs = Presentation(str(path))
    slides_data = []

    images_dir = ensure_dir(cache_dir / "slide_images")

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_info = {
            "source": "slides",
            "type": "pptx",
            "slide_number": slide_num,
            "title": "",
            "content": [],
            "notes": "",
            "images": []
        }

        # Extract title
        if slide.shapes.title:
            slide_info["title"] = slide.shapes.title.text

        # Extract text content
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text and text != slide_info["title"]:
                    slide_info["content"].append(text)

        # Extract notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                slide_info["notes"] = notes_slide.notes_text_frame.text

        # Extract images (basic - just note their presence)
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture type
                image_name = f"slide_{slide_num}_img_{len(slide_info['images'])}.png"
                slide_info["images"].append(image_name)

        slides_data.append(slide_info)

    logger.info(f"Parsed {len(slides_data)} slides from PPTX")
    return slides_data


def parse_pdf_slides(path: Path, cache_dir: Path, logger: logging.Logger) -> List[Dict[str, Any]]:
    """
    Parse a PDF as a sequence of slide-like dictionaries.
    
    When a page has few embedded characters, attempts OCR on a rendered high-resolution image; uses the first non-empty line of page text as the slide title and remaining lines as content. Extracts image references for each page into the `images` list.
    
    Returns:
        List of dictionaries, each with keys: `source`, `type`, `slide_number`, `title`, `content`, and `images`.
    """
    from examkit.ingestion.ocr import extract_text_with_ocr

    logger.info(f"Parsing PDF slides: {path}")

    doc = fitz.open(str(path))
    slides_data = []

    for page_num in range(len(doc)):
        page = doc[page_num]

        slide_info = {
            "source": "slides",
            "type": "pdf",
            "slide_number": page_num + 1,
            "title": "",
            "content": [],
            "images": []
        }

        # Extract text
        text = page.get_text()

        # Check text density - if low, might need OCR
        char_count = len(text.strip())
        if char_count < 50:
            logger.debug(f"Low text density on page {page_num + 1}, attempting OCR")
            # Render page to image for OCR
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom
            img_path = cache_dir / f"page_{page_num + 1}.png"
            pix.save(str(img_path))

            try:
                text = extract_text_with_ocr(img_path, logger)
            except Exception as e:
                logger.warning(f"OCR failed for page {page_num + 1}: {e}")

        # Split into lines and identify potential title
        lines = [line.strip() for line in text.split('\n') if line.strip()]
        if lines:
            # First substantial line is likely the title
            slide_info["title"] = lines[0]
            slide_info["content"] = lines[1:]

        # Extract images
        image_list = page.get_images()
        slide_info["images"] = [f"img_{img[0]}" for img in image_list]

        slides_data.append(slide_info)

    doc.close()
    logger.info(f"Parsed {len(slides_data)} slides from PDF")
    return slides_data


def parse_slides(path: Path, cache_dir: Path, logger: logging.Logger) -> List[Dict[str, Any]]:
    """
    Parse slides from PPTX or PDF format.

    Args:
        path: Path to slides file.
        cache_dir: Directory for cached files.
        logger: Logger instance.

    Returns:
        List of slide dictionaries.
    """
    suffix = path.suffix.lower()

    if suffix == '.pptx':
        return parse_pptx(path, cache_dir, logger)
    elif suffix == '.pdf':
        return parse_pdf_slides(path, cache_dir, logger)
    else:
        raise ValueError(f"Unsupported slides format: {suffix}")